from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import psutil
import win32com.client as win32
import subprocess
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Define pastel header colors (RGB)
PASTEL_COLORS = {
    "Motivation & Problem": RGBColor(255, 230, 230),   # soft pink
    "Approach & Pipeline":  RGBColor(230, 240, 255),   # soft blue
    "Results":              RGBColor(230, 255, 230),   # soft green
    "From Prototype to Production-Ready REM": RGBColor(255, 245, 225), # soft peach
    "Future Work":          RGBColor(240, 230, 255),   # soft lavender
}


OUTPUT_PPTX = "REMV_SPLASH_POSTER.pptx"
OUTPUT_PDF  = "REMV_SPLASH_POSTER.pdf"

def close_powerpoint():
    """Kill PowerPoint if running."""
    for proc in psutil.process_iter(['pid', 'name']):
        if proc.info['name'] and "POWERPNT.EXE" in proc.info['name']:
            print("Closing PowerPoint...")
            proc.kill()

def export_pdf(pptx_path, pdf_path):
    """Export a PPTX file to PDF using PowerPoint COM automation."""
    ppt_app = win32.Dispatch("PowerPoint.Application")
    ppt_app.Visible = 1
    presentation = ppt_app.Presentations.Open(os.path.abspath(pptx_path))
    presentation.SaveAs(os.path.abspath(pdf_path), 32)  # 32 = PDF format
    presentation.Close()
    ppt_app.Quit()
    print(f"Exported {pdf_path}")

def open_powerpoint(pptx_path):
    """Open file in PowerPoint."""
    subprocess.Popen(["powerpnt.exe", os.path.abspath(pptx_path)])
    print("Reopened PowerPoint.")

prs = Presentation()

# Define custom slide size: 1000 mm x 1000 mm (1m x 1m)
prs.slide_width = Mm(1000)
prs.slide_height = Mm(1000)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

def add_textbox(slide, left, top, width, height, text, font_size=32, bold=False, align=PP_ALIGN.LEFT, vcenter=False):
    txBox = slide.shapes.add_textbox(Mm(left), Mm(top), Mm(width), Mm(height))
    tf = txBox.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = align
    if vcenter:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return txBox

def add_section(slide, left, top, width, height,
                      header, body,
                      header_size=60, body_size=36,
                      box_color=RGBColor(240,240,240),
                      radius_mm=5):
    """
    Adds a section with a rounded header box and a rounded body box below,
    both sharing the same corner radius.
    """
    header_height = height * 0.2
    body_height   = height * 0.75

    # --- Header box ---
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Mm(left),
        Mm(top),
        Mm(width),
        Mm(header_height),
    )
    header_shape.adjustments[0] = radius_mm / (width / 2)  # consistent radius
    fill = header_shape.fill
    fill.solid()
    fill.fore_color.rgb = box_color
    header_shape.line.color.rgb = RGBColor(0,0,0)   # outline black
    header_shape.line.width = Pt(2)

    # Add header text (centered, bold, black)
    tf = header_shape.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = header
    p.font.size = Pt(header_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)  # force black text
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Body box ---
    body_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Mm(left),
        Mm(top + header_height + 5),   # little gap
        Mm(width),
        Mm(body_height),
    )
    body_shape.adjustments[0] = radius_mm / (width / 2)  # same radius
    fill = body_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255,255,255)     # white background
    body_shape.line.color.rgb = RGBColor(0,0,0)     # black outline
    body_shape.line.width = Pt(1.5)

    # Add body text (left aligned, black)
    tf = body_shape.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = body
    p.font.size = Pt(body_size)
    p.font.color.rgb = RGBColor(0, 0, 0)  # black
    p.alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.TOP

def add_image(slide, left, top, width, height, path="poster\\images\\sample.png"):
    """Add an image without a caption (labels are separate)."""
    return slide.shapes.add_picture(path, Mm(left), Mm(top), Mm(width), Mm(height))

# Column setup (3 columns top, title, 3 columns bottom)
col_width = 1000 / 3
margin = 20

# Top row (400 mm tall)
add_section(
    slide,
    margin,
    20,
    col_width - 2 * margin,
    400,
    "Motivation & Problem",
    "- Rust refactoring is hard (ownership, lifetimes, effects)\n"
    "- Compilation success ≠ semantic equivalence\n"
    "- Need trustworthy, automated refactoring",
    box_color=PASTEL_COLORS["Motivation & Problem"],
)

add_section(
    slide,
    col_width + margin,
    20,
    col_width - 2 * margin,
    400,
    "Approach & Pipeline",
    "[Diagram Placeholder]\n\nREM → CHARON → Aeneas → Coq\nZero-annotation, IDE integration",
    box_color=PASTEL_COLORS["Approach & Pipeline"],
)

add_section(
    slide,
    2 * col_width + margin,
    20,
    col_width - 2 * margin,
    400,
    "Results",
    "- 10/10 cases discharged\n- Avg cycle: 2s (IDE-friendly)\n\n[Table/Graph Placeholder]",
    box_color=PASTEL_COLORS["Results"],
)

# Middle (200 mm tall, title + authors + logos)
add_textbox(
    slide,
    20,
    420,
    960,
    100,
    "Verifying Extract Method Refactoring in Rust",
    font_size=100,
    bold=True,
    align=PP_ALIGN.CENTER,
    vcenter=True,
)

add_textbox(
    slide,
    20,
    500,
    960,
    60,
    "Matthew Britton, Alex Potanin, Sasha Pak",
    font_size=60,
    align=PP_ALIGN.CENTER,
    vcenter=True,
)

# Logos (in title band)
add_image(slide, 20, 430, 150, 100, path="poster/images/anu.png")    # ANU logo
add_image(slide, 865, 430, 120, 120*(349/429), path="poster/images/splash25.png")   # SPLASH logo

# Bottom row (400 mm tall)
add_section(
    slide,
    margin,
    580,
    col_width - 2 * margin,
    400,
    "From Prototype to Production-Ready REM",
    "- Standalone CLI\n- Async/await, generics, macros\n- VSCode extension",
    box_color=PASTEL_COLORS["From Prototype to Production-Ready REM"],
)

add_section(
    slide,
    col_width + margin,
    580,
    col_width - 2 * margin,
    400,
    "Future Work",
    "- Unsafe code\n- Concurrency\n- Large-scale evaluation\n- Better diagnostics\n- Scalable verification",
    box_color=PASTEL_COLORS["Future Work"],
)

# Right bottom column: QR codes styled in rounded boxes
qr_size = 90
qr_spacing = 130
card_width = 300
card_height = 115

def add_qr_card(slide, top, path, label, color):
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Mm(2 * col_width + margin),
        Mm(top),
        Mm(card_width),
        Mm(card_height),
    )
    # Fill background with chosen color
    fill = card.fill
    fill.solid()
    fill.fore_color.rgb = color
    card.line.fill.background()  # remove outline

    # Add QR code inside card (right side)
    add_image(
        slide,
        2 * col_width + margin + card_width - qr_size - 10,
        top + (card_height - qr_size) / 2,
        qr_size,
        qr_size,
        path=path,
    )

    # Add label inside card (left side, vertically centered, white text)
    tb = add_textbox(
        slide,
        2 * col_width + margin + 10,
        top,
        card_width - qr_size - 30,
        card_height,
        label,
        font_size=60,
        bold=True,
        align=PP_ALIGN.LEFT,
        vcenter=True,
    )
    # Force text color white
    for p in tb.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)

# Add QR code cards
add_qr_card(slide, 580, "poster\\images\\github.png", "GitHub", RGBColor(128, 0, 128))   # purple
add_qr_card(slide, 580 + qr_spacing, "poster\\images\\vscode.png", "VSCode", RGBColor(0, 102, 204))  # blue
add_qr_card(slide, 580 + 2 * qr_spacing, "poster\\images\\doi.png", "DOI", RGBColor(0, 0, 0))   # black

prs.save("REMV_SPLASH_POSTER.pptx")
print("Saved Poster")

def main():
    export_pdf(OUTPUT_PPTX, OUTPUT_PDF)
    # close_powerpoint()

if __name__ == "__main__":
    main()